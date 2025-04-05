#!/usr/bin/env python3
"""
Service implementation for PowerPoint operations.
"""
import os
import logging
import tempfile
import platform
import time
from datetime import datetime
from typing import Dict, List, Any, Optional, Union, Tuple
from pathlib import Path

from app.tools.base.service import ToolServiceBase


class PowerPointService(ToolServiceBase):
    """Service to handle PowerPoint operations across platforms"""

    def __init__(self):
        """Initialize the PowerPoint service"""
        super().__init__()
        self._app = None
        self._active_presentations = {}
        self._session_counter = 0
        self.temp_dir = self.get_env_var(
            "PPT_TEMP_DIR", default=tempfile.gettempdir())
        self.use_win32com = platform.system() == "Windows" and self.get_env_var(
            "PPT_USE_WIN32COM", default="True").lower() in ('true', 'yes', '1', 't', 'y')

    def validate(self) -> Tuple[bool, Optional[str]]:
        """
        Validate PowerPoint dependencies are available.

        Returns:
            Tuple[bool, str]: (is_valid, error_message)
        """
        try:
            if platform.system() == "Windows" and self.use_win32com:
                import win32com.client
            else:
                import pptx
            return True, None
        except ImportError as e:
            return False, f"Missing required dependencies: {str(e)}"

    def initialize(self) -> bool:
        """
        Initialize the PowerPoint service.

        Returns:
            bool: True if initialization was successful, False otherwise
        """
        try:
            # Verify dependencies
            is_valid, error_message = self.validate()
            if not is_valid:
                self.logger.error(
                    f"PowerPoint service validation failed: {error_message}")
                return False

            # Initialize based on platform
            if platform.system() == "Windows" and self.use_win32com:
                try:
                    import win32com.client
                    self._app = win32com.client.Dispatch(
                        "PowerPoint.Application")
                    self.logger.info(
                        "Initialized PowerPoint service using win32com")
                except Exception as e:
                    self.logger.error(
                        f"Failed to initialize PowerPoint with win32com: {str(e)}")
                    self.use_win32com = False

            # Ensure temp directory exists
            os.makedirs(self.temp_dir, exist_ok=True)

            self.initialized = True
            self.logger.info("PowerPoint service initialized successfully")
            return True
        except Exception as e:
            self.logger.error(
                f"Failed to initialize PowerPoint service: {str(e)}")
            return False

    def _get_next_session_id(self) -> str:
        """
        Generate a unique session ID.

        Returns:
            str: Unique session ID
        """
        self._session_counter += 1
        return f"ppt_session_{self._session_counter}_{int(time.time())}"

    def _cleanup_session(self, session_id: str) -> bool:
        """
        Clean up resources for a session.

        Args:
            session_id: ID of the session to clean up

        Returns:
            bool: True if cleanup was successful
        """
        if session_id not in self._active_presentations:
            return False

        if platform.system() == "Windows" and self.use_win32com:
            # Close the presentation in PowerPoint but don't quit the application
            try:
                presentation = self._active_presentations[session_id]["presentation"]
                if hasattr(presentation, "Close"):
                    presentation.Close()
            except Exception as e:
                self.logger.warning(f"Error closing presentation: {str(e)}")

        # Remove from active presentations
        del self._active_presentations[session_id]
        return True

    def create_presentation(self, session_id: str, template_path: Optional[str] = None) -> Dict[str, Any]:
        """
        Create a new PowerPoint presentation.

        Args:
            session_id: Unique identifier for the presentation session
            template_path: Optional path to a template file

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            if session_id in self._active_presentations:
                self._cleanup_session(session_id)

            session_data = {
                "presentation": None,
                "file_path": None,
                "created_at": datetime.now(),
                "modified_at": datetime.now(),
                "platform_type": "win32com" if (platform.system() == "Windows" and self.use_win32com) else "pptx"
            }

            if platform.system() == "Windows" and self.use_win32com:
                # Windows-specific implementation using win32com
                import win32com.client

                if not self._app:
                    self._app = win32com.client.Dispatch(
                        "PowerPoint.Application")

                if template_path and os.path.exists(template_path):
                    presentation = self._app.Presentations.Open(
                        os.path.abspath(template_path))
                else:
                    presentation = self._app.Presentations.Add()

                session_data["presentation"] = presentation

            else:
                # Cross-platform implementation using python-pptx
                from pptx import Presentation

                if template_path and os.path.exists(template_path):
                    presentation = Presentation(template_path)
                else:
                    presentation = Presentation()

                session_data["presentation"] = presentation

            self._active_presentations[session_id] = session_data

            return {
                "status": "success",
                "message": f"Created new presentation with session ID: {session_id}",
                "session_id": session_id,
                "platform_type": session_data["platform_type"]
            }

        except Exception as e:
            self.logger.error(f"Error creating presentation: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to create presentation: {str(e)}"
            }

    def open_presentation(self, session_id: str, file_path: str) -> Dict[str, Any]:
        """
        Open an existing PowerPoint presentation.

        Args:
            session_id: Unique identifier for the presentation session
            file_path: Path to the PowerPoint file

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            # Validate file path
            if not os.path.exists(file_path):
                return {
                    "status": "error",
                    "message": f"File not found: {file_path}"
                }

            if not file_path.lower().endswith(('.ppt', '.pptx')):
                return {
                    "status": "error",
                    "message": f"Not a PowerPoint file: {file_path}"
                }

            # Clean up existing session if needed
            if session_id in self._active_presentations:
                self._cleanup_session(session_id)

            session_data = {
                "presentation": None,
                "file_path": file_path,
                "created_at": datetime.now(),
                "modified_at": datetime.now(),
                "platform_type": "win32com" if (platform.system() == "Windows" and self.use_win32com) else "pptx"
            }

            if platform.system() == "Windows" and self.use_win32com:
                # Windows-specific implementation using win32com
                import win32com.client

                if not self._app:
                    self._app = win32com.client.Dispatch(
                        "PowerPoint.Application")

                presentation = self._app.Presentations.Open(
                    os.path.abspath(file_path))
                session_data["presentation"] = presentation

            else:
                # Cross-platform implementation using python-pptx
                from pptx import Presentation
                presentation = Presentation(file_path)
                session_data["presentation"] = presentation

            self._active_presentations[session_id] = session_data

            return {
                "status": "success",
                "message": f"Opened presentation from {file_path}",
                "session_id": session_id,
                "file_path": file_path,
                "platform_type": session_data["platform_type"]
            }

        except Exception as e:
            self.logger.error(f"Error opening presentation: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to open presentation: {str(e)}"
            }

    def save_presentation(self, session_id: str, file_path: Optional[str] = None) -> Dict[str, Any]:
        """
        Save the active PowerPoint presentation.

        Args:
            session_id: ID of the presentation session
            file_path: Optional path to save the file (if not provided, uses the existing path or generates a temporary one)

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            if session_id not in self._active_presentations:
                return {
                    "status": "error",
                    "message": f"Session not found: {session_id}"
                }

            session_data = self._active_presentations[session_id]
            presentation = session_data["presentation"]

            # Determine save path
            save_path = file_path or session_data["file_path"]
            if not save_path:
                save_path = os.path.join(self.temp_dir, f"{session_id}.pptx")

            # Ensure directory exists
            os.makedirs(os.path.dirname(
                os.path.abspath(save_path)), exist_ok=True)

            if session_data["platform_type"] == "win32com":
                # Save using win32com
                if hasattr(presentation, "SaveAs"):
                    presentation.SaveAs(os.path.abspath(save_path))
                else:
                    presentation.Save()
            else:
                # Save using python-pptx
                presentation.save(save_path)

            # Update session data
            session_data["file_path"] = save_path
            session_data["modified_at"] = datetime.now()

            return {
                "status": "success",
                "message": f"Saved presentation to {save_path}",
                "file_path": save_path
            }

        except Exception as e:
            self.logger.error(f"Error saving presentation: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to save presentation: {str(e)}"
            }

    def add_slide(self, session_id: str, layout_index: int = 1, title: Optional[str] = None,
                  content: Optional[str] = None) -> Dict[str, Any]:
        """
        Add a new slide to the presentation.

        Args:
            session_id: ID of the presentation session
            layout_index: Index of the slide layout to use
            title: Optional title for the slide
            content: Optional content for the slide

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            if session_id not in self._active_presentations:
                return {
                    "status": "error",
                    "message": f"Session not found: {session_id}"
                }

            session_data = self._active_presentations[session_id]
            presentation = session_data["presentation"]
            platform_type = session_data["platform_type"]

            if platform_type == "win32com":
                # Add slide using win32com
                slide_layouts = presentation.SlideMaster.CustomLayouts

                # Handle layout index bounds
                layout_count = slide_layouts.Count
                actual_layout_index = min(
                    layout_index, layout_count - 1) if layout_count > 0 else 0

                # Add the slide
                slide = presentation.Slides.AddSlide(
                    presentation.Slides.Count + 1,
                    # 1-indexed in win32com
                    slide_layouts.Item(actual_layout_index + 1)
                )

                # Add title if provided
                if title and hasattr(slide, "Shapes") and hasattr(slide.Shapes, "Title"):
                    slide.Shapes.Title.TextFrame.TextRange.Text = title

                # Add content if provided
                if content:
                    # Try to find a content placeholder
                    for shape_idx in range(1, slide.Shapes.Count + 1):
                        shape = slide.Shapes.Item(shape_idx)
                        # Content, Text, Body
                        if hasattr(shape, "PlaceholderFormat") and shape.PlaceholderFormat.Type in (2, 3, 7):
                            shape.TextFrame.TextRange.Text = content
                            break

                slide_index = slide.SlideIndex - 1  # Convert to 0-indexed

            else:
                # Add slide using python-pptx
                from pptx import Presentation

                # Handle layout index bounds
                if presentation.slide_layouts:
                    actual_layout_index = min(
                        layout_index, len(presentation.slide_layouts) - 1)
                    layout = presentation.slide_layouts[actual_layout_index]
                else:
                    # Fallback if no layouts available
                    layout = None

                # Add the slide
                if layout:
                    slide = presentation.slides.add_slide(layout)
                else:
                    # Basic slide without layout
                    slide = presentation.slides.add_slide(
                        presentation.slide_layouts[0])

                # Add title if provided
                if title and hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
                    slide.shapes.title.text = title

                # Add content if provided
                if content:
                    # Try to find a content placeholder
                    for shape in slide.placeholders:
                        # Content, Text, Body
                        if shape.placeholder_format.type in (2, 3, 7):
                            shape.text = content
                            break

                slide_index = len(presentation.slides) - 1

            # Update session data
            session_data["modified_at"] = datetime.now()

            return {
                "status": "success",
                "message": f"Added slide with layout index {layout_index}",
                "slide_index": slide_index
            }

        except Exception as e:
            self.logger.error(f"Error adding slide: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to add slide: {str(e)}"
            }

    def add_text(self, session_id: str, slide_index: int, text: str,
                 left: float = 1.0, top: float = 1.0, width: float = 8.0, height: float = 1.0,
                 font_size: int = 18, font_name: str = 'Calibri', bold: bool = False,
                 italic: bool = False, color: str = '000000') -> Dict[str, Any]:
        """
        Add text box to a slide.

        Args:
            session_id: ID of the presentation session
            slide_index: Index of the slide (0-indexed)
            text: Text content to add
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            font_size: Font size in points
            font_name: Font name
            bold: Whether the text should be bold
            italic: Whether the text should be italic
            color: Text color as hex string (without #)

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            if session_id not in self._active_presentations:
                return {
                    "status": "error",
                    "message": f"Session not found: {session_id}"
                }

            session_data = self._active_presentations[session_id]
            presentation = session_data["presentation"]
            platform_type = session_data["platform_type"]

            if platform_type == "win32com":
                # Add text using win32com
                # Get the slide (1-indexed in win32com)
                try:
                    slide = presentation.Slides.Item(slide_index + 1)
                except:
                    return {
                        "status": "error",
                        "message": f"Slide index {slide_index} out of range"
                    }

                # Convert inches to points (72 points per inch)
                left_points = left * 72
                top_points = top * 72
                width_points = width * 72
                height_points = height * 72

                # Add text box
                textbox = slide.Shapes.AddTextbox(
                    1,  # Orientation: msoTextOrientationHorizontal
                    left_points,
                    top_points,
                    width_points,
                    height_points
                )

                # Set text
                textframe = textbox.TextFrame
                textframe.TextRange.Text = text

                # Format text
                textrange = textframe.TextRange
                textrange.Font.Size = font_size
                textrange.Font.Name = font_name
                textrange.Font.Bold = bold
                textrange.Font.Italic = italic

                # Set color (expecting hex color code without #)
                if color and len(color) == 6:
                    try:
                        r = int(color[0:2], 16)
                        g = int(color[2:4], 16)
                        b = int(color[4:6], 16)
                        textrange.Font.Color.RGB = r + (g << 8) + (b << 16)
                    except ValueError:
                        pass  # Invalid color code, ignore

            else:
                # Add text using python-pptx
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor

                # Get the slide
                try:
                    slide = presentation.slides[slide_index]
                except IndexError:
                    return {
                        "status": "error",
                        "message": f"Slide index {slide_index} out of range"
                    }

                # Convert measurements to PowerPoint's units (inches)
                left_inches = Inches(left)
                top_inches = Inches(top)
                width_inches = Inches(width)
                height_inches = Inches(height)

                # Add text box
                textbox = slide.shapes.add_textbox(
                    left_inches, top_inches, width_inches, height_inches)
                textframe = textbox.text_frame
                textframe.text = text

                # Format text
                paragraph = textframe.paragraphs[0]
                run = paragraph.runs[0]
                run.font.size = Pt(font_size)
                run.font.name = font_name
                run.font.bold = bold
                run.font.italic = italic

                # Set color (expecting hex color code without #)
                if color and len(color) == 6:
                    try:
                        r = int(color[0:2], 16)
                        g = int(color[2:4], 16)
                        b = int(color[4:6], 16)
                        run.font.color.rgb = RGBColor(r, g, b)
                    except ValueError:
                        pass  # Invalid color code, ignore

            # Update session data
            session_data["modified_at"] = datetime.now()

            return {
                "status": "success",
                "message": f"Added text box to slide {slide_index}"
            }

        except Exception as e:
            self.logger.error(f"Error adding text: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to add text: {str(e)}"
            }

    def add_image(self, session_id: str, slide_index: int, image_path: str,
                  left: float = 1.0, top: float = 1.0, width: Optional[float] = None,
                  height: Optional[float] = None) -> Dict[str, Any]:
        """
        Add image to a slide.

        Args:
            session_id: ID of the presentation session
            slide_index: Index of the slide (0-indexed)
            image_path: Path to the image file
            left: Left position in inches
            top: Top position in inches
            width: Optional width in inches (if None, maintains aspect ratio)
            height: Optional height in inches (if None, maintains aspect ratio)

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            if session_id not in self._active_presentations:
                return {
                    "status": "error",
                    "message": f"Session not found: {session_id}"
                }

            # Check if image file exists
            if not os.path.exists(image_path):
                return {
                    "status": "error",
                    "message": f"Image file not found: {image_path}"
                }

            session_data = self._active_presentations[session_id]
            presentation = session_data["presentation"]
            platform_type = session_data["platform_type"]

            if platform_type == "win32com":
                # Add image using win32com
                # Get the slide (1-indexed in win32com)
                try:
                    slide = presentation.Slides.Item(slide_index + 1)
                except:
                    return {
                        "status": "error",
                        "message": f"Slide index {slide_index} out of range"
                    }

                # Convert inches to points (72 points per inch)
                left_points = left * 72
                top_points = top * 72

                # Add the image
                image = slide.Shapes.AddPicture(
                    os.path.abspath(image_path),
                    False,  # LinkToFile
                    True,   # SaveWithDocument
                    left_points,
                    top_points
                )

                # Resize if specified
                if width is not None and height is not None:
                    image.Width = width * 72
                    image.Height = height * 72
                elif width is not None:
                    # Maintain aspect ratio
                    aspect_ratio = image.Height / image.Width
                    image.Width = width * 72
                    image.Height = image.Width * aspect_ratio
                elif height is not None:
                    # Maintain aspect ratio
                    aspect_ratio = image.Width / image.Height
                    image.Height = height * 72
                    image.Width = image.Height * aspect_ratio

            else:
                # Add image using python-pptx
                from pptx import Presentation
                from pptx.util import Inches

                # Get the slide
                try:
                    slide = presentation.slides[slide_index]
                except IndexError:
                    return {
                        "status": "error",
                        "message": f"Slide index {slide_index} out of range"
                    }

                # Convert measurements to PowerPoint's units (inches)
                left_inches = Inches(left)
                top_inches = Inches(top)

                # Determine image size
                if width is not None and height is not None:
                    width_inches = Inches(width)
                    height_inches = Inches(height)
                    image = slide.shapes.add_picture(
                        image_path, left_inches, top_inches, width_inches, height_inches)
                else:
                    image = slide.shapes.add_picture(
                        image_path, left_inches, top_inches)

                    # Apply single dimension resize if specified
                    if width is not None:
                        aspect_ratio = image.height / image.width
                        image.width = Inches(width)
                        image.height = image.width * aspect_ratio
                    elif height is not None:
                        aspect_ratio = image.width / image.height
                        image.height = Inches(height)
                        image.width = image.height * aspect_ratio

            # Update session data
            session_data["modified_at"] = datetime.now()

            return {
                "status": "success",
                "message": f"Added image to slide {slide_index}"
            }

        except Exception as e:
            self.logger.error(f"Error adding image: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to add image: {str(e)}"
            }

    def add_chart(self, session_id: str, slide_index: int, chart_type: str,
                  categories: List[str], series_names: List[str], series_values: List[List[float]],
                  left: float = 1.0, top: float = 1.0, width: float = 8.0, height: float = 5.0,
                  chart_title: Optional[str] = None) -> Dict[str, Any]:
        """
        Add chart to a slide.

        Args:
            session_id: ID of the presentation session
            slide_index: Index of the slide (0-indexed)
            chart_type: Type of chart ('column', 'bar', 'line', 'pie', etc.)
            categories: List of category labels
            series_names: List of series names
            series_values: List of lists containing values for each series
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            chart_title: Optional title for the chart

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            if session_id not in self._active_presentations:
                return {
                    "status": "error",
                    "message": f"Session not found: {session_id}"
                }

            session_data = self._active_presentations[session_id]
            presentation = session_data["presentation"]
            platform_type = session_data["platform_type"]

            # Map chart type string to PowerPoint chart type
            chart_type_map = {
                'column': 3,    # Excel constant: xlColumnClustered
                'bar': 57,      # Excel constant: xlBarClustered
                'line': 4,      # Excel constant: xlLine
                'pie': 5,       # Excel constant: xlPie
                'area': 1,      # Excel constant: xlArea
                'scatter': -4169,  # Excel constant: xlXYScatter
                'doughnut': -4120,  # Excel constant: xlDoughnut
                'radar': -4151,  # Excel constant: xlRadar
                'bubble': 15    # Excel constant: xlBubble
            }

            if platform_type == "win32com":
                # Add chart using win32com
                # Get the slide (1-indexed in win32com)
                try:
                    slide = presentation.Slides.Item(slide_index + 1)
                except:
                    return {
                        "status": "error",
                        "message": f"Slide index {slide_index} out of range"
                    }

                # Convert inches to points (72 points per inch)
                left_points = left * 72
                top_points = top * 72
                width_points = width * 72
                height_points = height * 72

                # Get chart type
                excel_chart_type = chart_type_map.get(
                    chart_type.lower(), 3)  # Default to column

                # Add chart
                chart = slide.Shapes.AddChart2(
                    -1,  # Style (default)
                    excel_chart_type,
                    left_points,
                    top_points,
                    width_points,
                    height_points
                ).Chart

                # Access the chart data
                chart_data = chart.ChartData

                # Clear existing data
                chart_data.Workbook.Worksheets(1).Range("A:Z").Clear()

                # Add categories
                for i, category in enumerate(categories):
                    chart_data.Workbook.Worksheets(
                        1).Cells(1, i + 2).Value = category

                # Add series
                for i, series_name in enumerate(series_names):
                    chart_data.Workbook.Worksheets(1).Cells(
                        i + 2, 1).Value = series_name

                    # Add values for this series
                    if i < len(series_values):
                        for j, value in enumerate(series_values[i]):
                            if j < len(categories):
                                chart_data.Workbook.Worksheets(
                                    1).Cells(i + 2, j + 2).Value = value

                # Refresh the chart
                chart.Refresh()

                # Set chart title if provided
                if chart_title:
                    chart.HasTitle = True
                    chart.ChartTitle.Text = chart_title
                else:
                    chart.HasTitle = False

            else:
                # Add chart using python-pptx
                from pptx import Presentation
                from pptx.util import Inches
                from pptx.chart.data import CategoryChartData
                from pptx.enum.chart import XL_CHART_TYPE

                # Get the slide
                try:
                    slide = presentation.slides[slide_index]
                except IndexError:
                    return {
                        "status": "error",
                        "message": f"Slide index {slide_index} out of range"
                    }

                # Convert measurements to PowerPoint's units (inches)
                left_inches = Inches(left)
                top_inches = Inches(top)
                width_inches = Inches(width)
                height_inches = Inches(height)

                # Map chart type string to python-pptx chart type
                xl_chart_type_map = {
                    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                    'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                    'line': XL_CHART_TYPE.LINE,
                    'pie': XL_CHART_TYPE.PIE,
                    'area': XL_CHART_TYPE.AREA,
                    'scatter': XL_CHART_TYPE.XY_SCATTER,
                    'radar': XL_CHART_TYPE.RADAR,
                    'doughnut': XL_CHART_TYPE.DOUGHNUT,
                    'bubble': XL_CHART_TYPE.BUBBLE
                }

                xl_chart_type = xl_chart_type_map.get(
                    chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)

                # Create chart data
                chart_data = CategoryChartData()
                chart_data.categories = categories

                # Add series
                for i, series_name in enumerate(series_names):
                    if i < len(series_values):
                        chart_data.add_series(series_name, series_values[i])

                # Add chart to slide
                chart = slide.shapes.add_chart(
                    xl_chart_type, left_inches, top_inches,
                    width_inches, height_inches, chart_data).chart

                # Set chart title if provided
                if chart_title:
                    chart.has_title = True
                    chart.chart_title.text_frame.text = chart_title
                else:
                    chart.has_title = False

            # Update session data
            session_data["modified_at"] = datetime.now()

            return {
                "status": "success",
                "message": f"Added {chart_type} chart to slide {slide_index}"
            }

        except Exception as e:
            self.logger.error(f"Error adding chart: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to add chart: {str(e)}"
            }

    def add_table(self, session_id: str, slide_index: int, rows: int, cols: int,
                  data: List[List[str]], left: float = 1.0, top: float = 1.0,
                  width: float = 8.0, height: float = 5.0) -> Dict[str, Any]:
        """
        Add table to a slide.

        Args:
            session_id: ID of the presentation session
            slide_index: Index of the slide (0-indexed)
            rows: Number of rows
            cols: Number of columns
            data: 2D list of data to populate the table
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            if session_id not in self._active_presentations:
                return {
                    "status": "error",
                    "message": f"Session not found: {session_id}"
                }

            session_data = self._active_presentations[session_id]
            presentation = session_data["presentation"]
            platform_type = session_data["platform_type"]

            if platform_type == "win32com":
                # Add table using win32com
                # Get the slide (1-indexed in win32com)
                try:
                    slide = presentation.Slides.Item(slide_index + 1)
                except:
                    return {
                        "status": "error",
                        "message": f"Slide index {slide_index} out of range"
                    }

                # Convert inches to points (72 points per inch)
                left_points = left * 72
                top_points = top * 72
                width_points = width * 72
                height_points = height * 72

                # Add table shape
                table = slide.Shapes.AddTable(
                    rows, cols, left_points, top_points, width_points, height_points).Table

                # Populate table data
                for r in range(min(rows, len(data))):
                    for c in range(min(cols, len(data[r]) if r < len(data) else 0)):
                        table.Cell(
                            r + 1, c + 1).Shape.TextFrame.TextRange.Text = str(data[r][c])

            else:
                # Add table using python-pptx
                from pptx import Presentation
                from pptx.util import Inches

                # Get the slide
                try:
                    slide = presentation.slides[slide_index]
                except IndexError:
                    return {
                        "status": "error",
                        "message": f"Slide index {slide_index} out of range"
                    }

                # Convert measurements to PowerPoint's units (inches)
                left_inches = Inches(left)
                top_inches = Inches(top)
                width_inches = Inches(width)
                height_inches = Inches(height)

                # Add table
                table = slide.shapes.add_table(rows, cols, left_inches, top_inches,
                                               width_inches, height_inches).table

                # Populate table data
                for r in range(min(rows, len(data))):
                    for c in range(min(cols, len(data[r]) if r < len(data) else 0)):
                        cell = table.cell(r, c)
                        cell.text = str(data[r][c])

            # Update session data
            session_data["modified_at"] = datetime.now()

            return {
                "status": "success",
                "message": f"Added table with {rows} rows and {cols} columns to slide {slide_index}"
            }

        except Exception as e:
            self.logger.error(f"Error adding table: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to add table: {str(e)}"
            }

    def analyze_presentation(self, session_id: str) -> Dict[str, Any]:
        """
        Analyze the content and structure of a presentation.

        Args:
            session_id: ID of the presentation session

        Returns:
            Dict with analysis results
        """
        self._is_initialized()
        try:
            if session_id not in self._active_presentations:
                return {
                    "status": "error",
                    "message": f"Session not found: {session_id}"
                }

            session_data = self._active_presentations[session_id]
            presentation = session_data["presentation"]
            platform_type = session_data["platform_type"]

            analysis = {
                "total_slides": 0,
                "slides": [],
                "word_count": 0,
                "total_images": 0,
                "total_charts": 0,
                "total_tables": 0,
                "average_words_per_slide": 0,
                "slide_titles": [],
                "presentation_structure": []
            }

            if platform_type == "win32com":
                # Analyze using win32com
                slides_count = presentation.Slides.Count
                analysis["total_slides"] = slides_count

                for i in range(1, slides_count + 1):
                    slide = presentation.Slides.Item(i)
                    slide_analysis = self._analyze_slide_win32com(slide)
                    slide_analysis["slide_number"] = i - \
                        1  # Convert to 0-indexed
                    analysis["slides"].append(slide_analysis)

                    # Update presentation-wide statistics
                    analysis["word_count"] += slide_analysis["word_count"]
                    analysis["total_images"] += slide_analysis["image_count"]
                    analysis["total_charts"] += slide_analysis["chart_count"]
                    analysis["total_tables"] += slide_analysis["table_count"]

                    # Get slide title
                    if slide_analysis["title"]:
                        analysis["slide_titles"].append(
                            slide_analysis["title"])
                        analysis["presentation_structure"].append({
                            "slide_number": i - 1,  # Convert to 0-indexed
                            "title": slide_analysis["title"]
                        })
                    else:
                        analysis["presentation_structure"].append({
                            "slide_number": i - 1,  # Convert to 0-indexed
                            "title": f"Slide {i} (No Title)"
                        })

            else:
                # Analyze using python-pptx
                slides_count = len(presentation.slides)
                analysis["total_slides"] = slides_count

                for i, slide in enumerate(presentation.slides):
                    slide_analysis = self._analyze_slide_pptx(slide)
                    slide_analysis["slide_number"] = i
                    analysis["slides"].append(slide_analysis)

                    # Update presentation-wide statistics
                    analysis["word_count"] += slide_analysis["word_count"]
                    analysis["total_images"] += slide_analysis["image_count"]
                    analysis["total_charts"] += slide_analysis["chart_count"]
                    analysis["total_tables"] += slide_analysis["table_count"]

                    # Get slide title
                    if slide_analysis["title"]:
                        analysis["slide_titles"].append(
                            slide_analysis["title"])
                        analysis["presentation_structure"].append({
                            "slide_number": i,
                            "title": slide_analysis["title"]
                        })
                    else:
                        analysis["presentation_structure"].append({
                            "slide_number": i,
                            "title": f"Slide {i+1} (No Title)"
                        })

            # Calculate averages
            if analysis["total_slides"] > 0:
                analysis["average_words_per_slide"] = analysis["word_count"] / \
                    analysis["total_slides"]

            return {
                "status": "success",
                "analysis": analysis
            }

        except Exception as e:
            self.logger.error(f"Error analyzing presentation: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to analyze presentation: {str(e)}"
            }

    def _analyze_slide_win32com(self, slide) -> Dict[str, Any]:
        """
        Analyze a single slide using win32com.

        Args:
            slide: Slide object from win32com

        Returns:
            Dict with slide analysis
        """
        slide_analysis = {
            "title": "",
            "word_count": 0,
            "text_content": "",
            "image_count": 0,
            "chart_count": 0,
            "table_count": 0,
            "elements": []
        }

        # Get slide title
        try:
            if hasattr(slide, 'Shapes') and hasattr(slide.Shapes, 'Title') and slide.Shapes.Title:
                slide_analysis["title"] = slide.Shapes.Title.TextFrame.TextRange.Text
        except:
            pass

        # Analyze shapes
        for i in range(1, slide.Shapes.Count + 1):
            shape = slide.Shapes.Item(i)

            # Text analysis
            if hasattr(shape, 'TextFrame') and shape.TextFrame.HasText:
                text = shape.TextFrame.TextRange.Text

                if text.strip():
                    slide_analysis["text_content"] += text + "\n"
                    words = text.split()
                    slide_analysis["word_count"] += len(words)

                    slide_analysis["elements"].append({
                        "type": "text",
                        "content": text.strip(),
                        "word_count": len(words)
                    })

            # Image analysis
            if shape.Type == 13:  # msoPicture
                slide_analysis["image_count"] += 1
                slide_analysis["elements"].append({
                    "type": "image"
                })

            # Chart analysis
            if hasattr(shape, 'HasChart') and shape.HasChart:
                slide_analysis["chart_count"] += 1
                slide_analysis["elements"].append({
                    "type": "chart",
                    "chart_type": str(shape.Chart.ChartType)
                })

            # Table analysis
            if hasattr(shape, 'HasTable') and shape.HasTable:
                slide_analysis["table_count"] += 1

                table = shape.Table
                rows = table.Rows.Count
                cols = table.Columns.Count

                table_data = []
                for r in range(1, rows + 1):
                    row_data = []
                    for c in range(1, cols + 1):
                        try:
                            cell_text = table.Cell(
                                r, c).Shape.TextFrame.TextRange.Text
                            row_data.append(cell_text)
                        except:
                            row_data.append("")
                    table_data.append(row_data)

                slide_analysis["elements"].append({
                    "type": "table",
                    "rows": rows,
                    "columns": cols,
                    "data": table_data
                })

        return slide_analysis

    def _analyze_slide_pptx(self, slide) -> Dict[str, Any]:
        """
        Analyze a single slide using python-pptx.

        Args:
            slide: Slide object from python-pptx

        Returns:
            Dict with slide analysis
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        slide_analysis = {
            "title": "",
            "word_count": 0,
            "text_content": "",
            "image_count": 0,
            "chart_count": 0,
            "table_count": 0,
            "elements": []
        }

        # Get slide title
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide_analysis["title"] = slide.shapes.title.text

        # Analyze shapes
        for shape in slide.shapes:
            # Text analysis
            if shape.has_text_frame:
                text = ""
                for paragraph in shape.text_frame.paragraphs:
                    text += paragraph.text + "\n"

                if text.strip():
                    slide_analysis["text_content"] += text
                    words = text.split()
                    slide_analysis["word_count"] += len(words)

                    slide_analysis["elements"].append({
                        "type": "text",
                        "content": text.strip(),
                        "word_count": len(words)
                    })

            # Image analysis
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_analysis["image_count"] += 1
                slide_analysis["elements"].append({
                    "type": "image"
                })

            # Chart analysis
            if shape.has_chart:
                slide_analysis["chart_count"] += 1
                slide_analysis["elements"].append({
                    "type": "chart",
                    "chart_type": str(shape.chart.chart_type)
                })

            # Table analysis
            if shape.has_table:
                slide_analysis["table_count"] += 1
                rows = len(shape.table.rows)
                cols = len(shape.table.columns)

                table_data = []
                for r in range(rows):
                    row_data = []
                    for c in range(cols):
                        cell_text = shape.table.cell(r, c).text
                        row_data.append(cell_text)
                    table_data.append(row_data)

                slide_analysis["elements"].append({
                    "type": "table",
                    "rows": rows,
                    "columns": cols,
                    "data": table_data
                })

        return slide_analysis

    def enhance_presentation(self, session_id: str) -> Dict[str, Any]:
        """
        Provide suggestions to enhance the presentation.

        Args:
            session_id: ID of the presentation session

        Returns:
            Dict with enhancement suggestions
        """
        self._is_initialized()
        try:
            # First analyze the presentation
            analysis_result = self.analyze_presentation(session_id)

            if analysis_result["status"] != "success":
                return analysis_result

            analysis = analysis_result["analysis"]

            suggestions = {
                "overall_suggestions": [],
                "slide_suggestions": []
            }

            # Overall suggestions
            self._add_overall_suggestions(analysis, suggestions)

            # Per-slide suggestions
            for slide_analysis in analysis["slides"]:
                slide_suggestions = self._generate_slide_suggestions(
                    slide_analysis)
                suggestions["slide_suggestions"].append({
                    "slide_number": slide_analysis["slide_number"],
                    "title": slide_analysis["title"] or f"Slide {slide_analysis['slide_number']+1}",
                    "suggestions": slide_suggestions
                })

            return {
                "status": "success",
                "suggestions": suggestions
            }

        except Exception as e:
            self.logger.error(
                f"Error generating enhancement suggestions: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to generate enhancement suggestions: {str(e)}"
            }

    def _add_overall_suggestions(self, analysis, suggestions):
        """
        Add overall presentation suggestions.

        Args:
            analysis: Presentation analysis data
            suggestions: Suggestions object to update
        """
        # Check for consistent structure
        if analysis["total_slides"] > 1:
            suggestions["overall_suggestions"].append({
                "type": "structure",
                "suggestion": "Consider adding an agenda or table of contents slide at the beginning"
            })

        # Check word count
        if analysis["average_words_per_slide"] > 100:
            suggestions["overall_suggestions"].append({
                "type": "content_density",
                "suggestion": "Presentation slides have too much text (average over 100 words per slide). Consider breaking content into more slides or reducing text."
            })

        # Check slide count
        if analysis["total_slides"] > 15:
            suggestions["overall_suggestions"].append({
                "type": "length",
                "suggestion": "Presentation is quite long. Consider condensing or breaking into multiple presentations."
            })

        # Check visual elements
        if analysis["total_images"] + analysis["total_charts"] < analysis["total_slides"] / 2:
            suggestions["overall_suggestions"].append({
                "type": "visuals",
                "suggestion": "Add more visual elements (images, charts) to increase engagement."
            })

    def _generate_slide_suggestions(self, slide_analysis):
        """
        Generate suggestions for a single slide.

        Args:
            slide_analysis: Analysis data for a slide

        Returns:
            List of suggestions for the slide
        """
        suggestions = []

        # Check for title
        if not slide_analysis["title"]:
            suggestions.append({
                "type": "structure",
                "suggestion": "Add a clear title to this slide"
            })

        # Check word count on slide
        if slide_analysis["word_count"] > 100:
            suggestions.append({
                "type": "content_density",
                "suggestion": "Slide has too much text. Consider breaking into multiple slides or reducing text."
            })
        elif slide_analysis["word_count"] < 10 and slide_analysis["image_count"] == 0 and slide_analysis["chart_count"] == 0:
            suggestions.append({
                "type": "content",
                "suggestion": "Slide has very little content. Consider adding more information or visuals."
            })

        # Check balance of elements
        if slide_analysis["word_count"] > 50 and slide_analysis["image_count"] + slide_analysis["chart_count"] == 0:
            suggestions.append({
                "type": "balance",
                "suggestion": "Text-heavy slide. Consider adding relevant images or charts to illustrate points."
            })

        return suggestions

    def generate_presentation(self, session_id: str, title: str, content: str) -> Dict[str, Any]:
        """
        Generate a presentation from text content.

        Args:
            session_id: ID of the presentation session
            title: Title for the presentation
            content: Text content to convert into a presentation

        Returns:
            Dict with result information
        """
        self._is_initialized()
        try:
            # Create new presentation
            result = self.create_presentation(session_id)

            if result["status"] != "success":
                return result

            # Process content to extract slide structure
            slides = self._extract_slide_structure(content)

            # Add title slide
            # Use first layout (usually title slide)
            self.add_slide(session_id, 0, title)

            # Add content slides
            slide_count = 1  # Start after title slide

            for slide_data in slides:
                slide_title = slide_data.get("title", f"Slide {slide_count}")
                slide_content = slide_data.get("content", "")

                # Use second layout (usually title and content)
                result = self.add_slide(
                    session_id, 1, slide_title, slide_content)

                if result["status"] != "success":
                    return {
                        "status": "error",
                        "message": f"Failed to add slide {slide_count}: {result['message']}"
                    }

                slide_count += 1

            return {
                "status": "success",
                "message": f"Generated presentation with {slide_count} slides",
                "slides_count": slide_count
            }

        except Exception as e:
            self.logger.error(f"Error generating presentation: {str(e)}")
            return {
                "status": "error",
                "message": f"Failed to generate presentation: {str(e)}"
            }

    def _extract_slide_structure(self, text: str) -> List[Dict[str, Any]]:
        """
        Extract slide structure from plain text content.

        Args:
            text: Raw text content

        Returns:
            List of slide data dictionaries
        """
        # Split text into paragraphs
        paragraphs = text.split('\n\n')

        # Extract title (first non-empty paragraph)
        title = "New Presentation"
        content_start = 0

        for i, paragraph in enumerate(paragraphs):
            if paragraph.strip():
                title = paragraph.strip()
                content_start = i + 1
                break

        # Initialize slides
        slides = []

        # Process remaining paragraphs for content slides
        current_slide = None

        for paragraph in paragraphs[content_start:]:
            paragraph = paragraph.strip()
            if not paragraph:
                continue

            # Check if this is a heading (potential slide title)
            sentences = paragraph.split('. ')
            first_sentence = sentences[0] if sentences else ""

            is_heading = (
                len(first_sentence) < 100 and
                len(sentences) <= 2 and
                not first_sentence.endswith('.')
            )

            if is_heading:
                # Save previous slide if exists
                if current_slide:
                    slides.append(current_slide)

                # Start new slide
                current_slide = {
                    "title": paragraph,
                    "content": ""
                }
            else:
                # Add content to current slide or create new one if needed
                if not current_slide:
                    current_slide = {
                        "title": "Content Slide",
                        "content": paragraph
                    }
                else:
                    current_slide["content"] += paragraph + "\n\n"

        # Add the last slide if not empty
        if current_slide:
            slides.append(current_slide)

        return slides


# Singleton instance
_service_instance = None


def get_service() -> PowerPointService:
    """
    Get or initialize the PowerPoint service singleton.

    Returns:
        PowerPointService instance
    """
    global _service_instance
    if _service_instance is None:
        _service_instance = PowerPointService()
        _service_instance.initialize()
    return _service_instance
